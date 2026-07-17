"""HAWK File Studio — belge ÜRET / DÜZENLE / DÖNÜŞTÜR + indirilebilir çıktı.
Saf yapı: LLM çağrısı YOK (app.py spec'i üretir, burası dosyayı kurar). Döner: bytes veya kaydedilmiş URL.

Desteklenen üretim: docx, xlsx (grafik dahil), pptx, pdf, csv, kod/metin.
Desteklenen dönüştürme: pdf/docx/xlsx/csv/txt arası (metin tabanlı) + zip listeleme/çıkarma.
"""
import io
import os
import csv as _csv
import uuid
import zipfile

UPLOAD_DIR = "/app/uploads/files"
PUBLIC_BASE = os.getenv("HAWK_PUBLIC_BASE", "https://www.hawk-operasyon.com")

_EXT_MIME = {
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "pdf": "application/pdf", "csv": "text/csv", "txt": "text/plain",
}


def save_file(raw: bytes, ext: str, prefix: str = "hawk") -> dict:
    """Bytes'ı uploads/files altına kaydet, indirme URL'i döndür."""
    os.makedirs(UPLOAD_DIR, exist_ok=True)
    ext = (ext or "bin").lstrip(".").lower()
    fn = f"{prefix}_{uuid.uuid4().hex[:10]}.{ext}"
    with open(os.path.join(UPLOAD_DIR, fn), "wb") as f:
        f.write(raw)
    # F4: kalıcı-anonim /api/dl kapalı → süreli, dosya-kapsamlı İMZALI URL döndür.
    try:
        from core.dl_sign import signed_url as _signed_url
        url = _signed_url(fn)
    except Exception:
        url = f"{PUBLIC_BASE}/api/dl/{fn}"
    return {"filename": fn, "url": url,
            "bytes": len(raw), "mime": _EXT_MIME.get(ext, "application/octet-stream")}


# ======================= ÜRETİCİLER (spec dict → bytes) =======================

def build_docx(spec: dict) -> bytes:
    from docx import Document
    from docx.shared import Pt
    doc = Document()
    if spec.get("title"):
        doc.add_heading(str(spec["title"])[:250], level=0)
    if spec.get("subtitle"):
        p = doc.add_paragraph(str(spec["subtitle"])); p.runs and setattr(p.runs[0].font, "size", Pt(12))
    for sec in spec.get("sections", []) or []:
        if not isinstance(sec, dict):
            doc.add_paragraph(str(sec)); continue
        if sec.get("heading"):
            doc.add_heading(str(sec["heading"])[:250], level=1)
        body = sec.get("body") or sec.get("paragraphs") or []
        if isinstance(body, str):
            body = [body]
        for para in body:
            doc.add_paragraph(str(para))
        for b in (sec.get("bullets") or []):
            doc.add_paragraph(str(b), style="List Bullet")
        # basit tablo
        tbl = sec.get("table")
        if isinstance(tbl, dict) and tbl.get("rows"):
            headers = tbl.get("headers") or []
            rows = tbl.get("rows") or []
            ncol = len(headers) or (len(rows[0]) if rows else 0)
            if ncol:
                t = doc.add_table(rows=0, cols=ncol); t.style = "Light Grid Accent 1"
                if headers:
                    hc = t.add_row().cells
                    for i, h in enumerate(headers[:ncol]):
                        hc[i].text = str(h)
                for r in rows:
                    rc = t.add_row().cells
                    for i, v in enumerate(list(r)[:ncol]):
                        rc[i].text = str(v)
    bio = io.BytesIO(); doc.save(bio); return bio.getvalue()


def build_xlsx(spec: dict) -> bytes:
    from openpyxl import Workbook
    from openpyxl.chart import BarChart, LineChart, PieChart, Reference
    wb = Workbook()
    sheets = spec.get("sheets")
    if not sheets:
        sheets = [{"name": spec.get("name", "Sheet1"),
                   "headers": spec.get("headers", []), "rows": spec.get("rows", []),
                   "chart": spec.get("chart")}]
    first = True
    for sh in sheets:
        ws = wb.active if first else wb.create_sheet()
        first = False
        ws.title = (str(sh.get("name") or "Sheet"))[:31]
        headers = sh.get("headers") or []
        rows = sh.get("rows") or []
        if headers:
            ws.append([str(h) for h in headers])
        for r in rows:
            ws.append(list(r))
        ch = sh.get("chart")
        if ch and rows and headers and len(headers) >= 2:
            try:
                ctype = str(ch.get("type", "bar")).lower()
                chart = {"line": LineChart, "pie": PieChart}.get(ctype, BarChart)()
                chart.title = ch.get("title") or ws.title
                nrows = len(rows)
                val_cols = ch.get("value_cols") or list(range(1, len(headers)))
                for vc in val_cols:
                    data = Reference(ws, min_col=vc + 1, min_row=1, max_row=nrows + 1)
                    chart.add_data(data, titles_from_data=True)
                cats = Reference(ws, min_col=1, min_row=2, max_row=nrows + 1)
                chart.set_categories(cats)
                anchor_col = chr(ord("A") + len(headers) + 1)
                ws.add_chart(chart, f"{anchor_col}2")
            except Exception:
                pass
    bio = io.BytesIO(); wb.save(bio); return bio.getvalue()


def build_pptx(spec: dict) -> bytes:
    from pptx import Presentation
    from pptx.util import Pt
    prs = Presentation()
    slides = spec.get("slides") or []
    if spec.get("title") and not slides:
        slides = [{"title": spec["title"], "bullets": []}]
    for i, sl in enumerate(slides):
        if not isinstance(sl, dict):
            sl = {"title": str(sl), "bullets": []}
        layout = prs.slide_layouts[0 if i == 0 else 1]
        slide = prs.slides.add_slide(layout)
        try:
            slide.shapes.title.text = str(sl.get("title", ""))[:200]
        except Exception:
            pass
        bullets = sl.get("bullets") or sl.get("body") or []
        if isinstance(bullets, str):
            bullets = [bullets]
        body_ph = None
        for ph in slide.placeholders:
            if ph.placeholder_format.idx != 0:
                body_ph = ph; break
        if body_ph is not None and bullets:
            tf = body_ph.text_frame
            tf.text = str(bullets[0])
            for b in bullets[1:]:
                p = tf.add_paragraph(); p.text = str(b)
    bio = io.BytesIO(); prs.save(bio); return bio.getvalue()


_PDF_FONT = None


def _pdf_font():
    """Türkçe karakter destekli TTF kaydet (DejaVuSans) — yoksa Helvetica."""
    global _PDF_FONT
    if _PDF_FONT is not None:
        return _PDF_FONT
    try:
        from reportlab.pdfbase import pdfmetrics
        from reportlab.pdfbase.ttfonts import TTFont
        for path in ["/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
                     "/usr/share/fonts/dejavu/DejaVuSans.ttf",
                     "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf"]:
            if os.path.exists(path):
                pdfmetrics.registerFont(TTFont("HawkUni", path))
                _PDF_FONT = "HawkUni"
                return _PDF_FONT
    except Exception:
        pass
    _PDF_FONT = "Helvetica"
    return _PDF_FONT


def build_pdf(spec: dict) -> bytes:
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.lib.units import cm
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, ListFlowable, ListItem
    font = _pdf_font()
    styles = getSampleStyleSheet()
    for s in styles.byName.values():
        try:
            s.fontName = font if "Bold" not in getattr(s, "fontName", "") else font
        except Exception:
            pass
    bio = io.BytesIO()
    doc = SimpleDocTemplate(bio, pagesize=A4, topMargin=2 * cm, bottomMargin=2 * cm,
                            leftMargin=2 * cm, rightMargin=2 * cm)
    story = []
    if spec.get("title"):
        story.append(Paragraph(str(spec["title"]), styles["Title"])); story.append(Spacer(1, 14))
    for sec in spec.get("sections", []) or []:
        if not isinstance(sec, dict):
            story.append(Paragraph(str(sec), styles["BodyText"])); continue
        if sec.get("heading"):
            story.append(Paragraph(str(sec["heading"]), styles["Heading2"]))
        body = sec.get("body") or sec.get("paragraphs") or []
        if isinstance(body, str):
            body = [body]
        for para in body:
            story.append(Paragraph(str(para), styles["BodyText"])); story.append(Spacer(1, 6))
        bullets = sec.get("bullets") or []
        if bullets:
            story.append(ListFlowable([ListItem(Paragraph(str(b), styles["BodyText"])) for b in bullets],
                                      bulletType="bullet"))
            story.append(Spacer(1, 6))
    if not story:
        story.append(Paragraph(str(spec.get("text", "")), styles["BodyText"]))
    doc.build(story)
    return bio.getvalue()


def build_csv(spec: dict) -> bytes:
    bio = io.StringIO()
    w = _csv.writer(bio)
    headers = spec.get("headers") or []
    if headers:
        w.writerow(headers)
    for r in spec.get("rows", []) or []:
        w.writerow(list(r))
    return bio.getvalue().encode("utf-8-sig")


def build_code(spec: dict) -> bytes:
    return str(spec.get("code", "")).encode("utf-8")


_BUILDERS = {"docx": build_docx, "xlsx": build_xlsx, "pptx": build_pptx,
             "pdf": build_pdf, "csv": build_csv}


def build(target_ext: str, spec: dict) -> bytes:
    fn = _BUILDERS.get((target_ext or "").lower())
    if not fn:
        return build_code(spec if isinstance(spec, dict) else {"code": str(spec)})
    return fn(spec if isinstance(spec, dict) else {})


# ======================= ÇIKARMA / DÖNÜŞTÜRME =======================

def extract_text(data: bytes, filename: str = "") -> str:
    """Kaynak dosyadan düz metin çıkar (pdf/docx/xlsx/csv/txt/kod)."""
    ext = (filename.rsplit(".", 1)[-1] if "." in filename else "").lower()
    try:
        if ext == "pdf":
            try:
                import fitz
                d = fitz.open(stream=data, filetype="pdf")
                return "\n".join(p.get_text() for p in d)
            except Exception:
                from pypdf import PdfReader
                r = PdfReader(io.BytesIO(data))
                return "\n".join((pg.extract_text() or "") for pg in r.pages)
        if ext == "docx":
            from docx import Document
            doc = Document(io.BytesIO(data))
            return "\n".join(p.text for p in doc.paragraphs)
        if ext in ("xlsx", "xlsm"):
            from openpyxl import load_workbook
            wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
            out = []
            for ws in wb.worksheets:
                out.append(f"# {ws.title}")
                for row in ws.iter_rows(values_only=True):
                    out.append("\t".join("" if c is None else str(c) for c in row))
            return "\n".join(out)
        if ext == "csv":
            return data.decode("utf-8", "ignore")
        return data.decode("utf-8", "ignore")
    except Exception as e:
        return f"[extract_error:{str(e)[:120]}]"


def read_tabular(data: bytes, filename: str = ""):
    """csv/xlsx → (headers, rows). Grafik/dönüştürme için."""
    ext = (filename.rsplit(".", 1)[-1] if "." in filename else "").lower()
    if ext in ("xlsx", "xlsm"):
        from openpyxl import load_workbook
        wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        ws = wb.active
        rows = [list(r) for r in ws.iter_rows(values_only=True)]
    else:
        text = data.decode("utf-8", "ignore")
        rows = [r for r in _csv.reader(io.StringIO(text))]
    if not rows:
        return [], []
    return [str(h) for h in rows[0]], rows[1:]


def zip_list(data: bytes):
    with zipfile.ZipFile(io.BytesIO(data)) as z:
        return [{"name": i.filename, "size": i.file_size} for i in z.infolist() if not i.is_dir()]


def zip_extract_text(data: bytes, max_files: int = 40) -> str:
    out = []
    with zipfile.ZipFile(io.BytesIO(data)) as z:
        for i in z.infolist()[:max_files]:
            if i.is_dir():
                continue
            out.append(f"\n===== {i.filename} ({i.file_size} bayt) =====")
            try:
                out.append(extract_text(z.read(i.filename), i.filename)[:4000])
            except Exception as e:
                out.append(f"[okunamadı: {str(e)[:80]}]")
    return "\n".join(out)


def build_zip(files: dict) -> bytes:
    """{yol: içerik(str|bytes)} → zip bytes. Kod projesi/çoklu dosya paketleme için."""
    bio = io.BytesIO()
    with zipfile.ZipFile(bio, "w", zipfile.ZIP_DEFLATED) as z:
        for path, content in (files or {}).items():
            safe = str(path).lstrip("/").replace("..", "_")
            data = content.encode("utf-8") if isinstance(content, str) else bytes(content or b"")
            z.writestr(safe, data)
    return bio.getvalue()


def convert(data: bytes, src_name: str, target_ext: str) -> bytes:
    """Metin/tablo tabanlı dönüştürme. Döner: bytes (hedef format)."""
    target_ext = (target_ext or "").lstrip(".").lower()
    src_ext = (src_name.rsplit(".", 1)[-1] if "." in src_name else "").lower()
    # tablo ↔ tablo
    if src_ext in ("csv", "xlsx", "xlsm") and target_ext in ("csv", "xlsx"):
        headers, rows = read_tabular(data, src_name)
        spec = {"headers": headers, "rows": rows}
        return build_xlsx(spec) if target_ext == "xlsx" else build_csv(spec)
    # metin çıkar → hedef belge
    text = extract_text(data, src_name)
    paras = [p for p in text.split("\n") if p.strip()]
    spec = {"title": src_name.rsplit(".", 1)[0],
            "sections": [{"body": paras}]}
    if target_ext in ("docx", "pdf"):
        return build(target_ext, spec)
    if target_ext == "txt":
        return text.encode("utf-8")
    if target_ext == "csv":
        return build_csv({"rows": [[p] for p in paras]})
    return text.encode("utf-8")
